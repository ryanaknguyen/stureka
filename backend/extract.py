import fitz  # PyMuPDF
import io
from pptx import Presentation  # python-pptx

def extract_text(filename, file_bytes):
    if filename.endswith(".pdf"):
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    elif filename.endswith(".txt"):
        return file_bytes.decode("utf-8")
    elif filename.endswith(".pptx"):
        presentation = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return "Unsupported file type"
